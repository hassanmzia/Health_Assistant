"""
Generate Technical Architecture PowerPoint Deck
Healthcare Intelligence Platform
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)
BG_CARD_LIGHT = RGBColor(0x27, 0x34, 0x49)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
BLUE = RGBColor(0x38, 0x82, 0xF6)
BLUE_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
BLUE_BG = RGBColor(0x1E, 0x3A, 0x5F)
GREEN = RGBColor(0x10, 0xB9, 0x81)
GREEN_BG = RGBColor(0x06, 0x4E, 0x3B)
PINK = RGBColor(0xEC, 0x48, 0x99)
PINK_BG = RGBColor(0x5B, 0x21, 0x46)
YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
YELLOW_BG = RGBColor(0x5C, 0x3D, 0x0C)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_BG = RGBColor(0x3B, 0x24, 0x6B)
RED = RGBColor(0xEF, 0x44, 0x44)
RED_BG = RGBColor(0x5C, 0x1A, 0x1A)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
INDIGO_BG = RGBColor(0x27, 0x28, 0x5B)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
TEAL_BG = RGBColor(0x0D, 0x4F, 0x4B)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1), radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_text(slide, left, top, width, height, lines, font_size=12, color=WHITE, font_name='Calibri', line_spacing=1.2):
    """lines = list of (text, bold, color, font_size_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        text = line_info[0]
        is_bold = line_info[1] if len(line_info) > 1 else False
        line_color = line_info[2] if len(line_info) > 2 else color
        fs = line_info[3] if len(line_info) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.color.rgb = line_color
        run.font.bold = is_bold
        run.font.name = font_name
    return txBox

def add_arrow(slide, start_left, start_top, end_left, end_top, color=GRAY_400, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_service_card(slide, left, top, width, height, title, subtitle, color, bg_color, items=None):
    add_shape(slide, left, top, width, height, bg_color, color, Pt(2))
    add_text(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.3), title, font_size=13, color=color, bold=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.38), width - Inches(0.4), Inches(0.25), subtitle, font_size=9, color=GRAY_400)
    if items:
        lines = [(item, False, GRAY_300, 9) for item in items]
        add_rich_text(slide, left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.7), lines, font_size=9, color=GRAY_300)


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide1, BG_DARK)

# Accent line
add_shape(slide1, Inches(2), Inches(2.5), Inches(9.333), Pt(3), BLUE)

add_text(slide1, Inches(2), Inches(2.7), Inches(9.333), Inches(1),
         "Healthcare Intelligence Platform", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide1, Inches(2), Inches(3.6), Inches(9.333), Inches(0.6),
         "Technical Architecture", font_size=28, color=BLUE_LIGHT, bold=False, alignment=PP_ALIGN.CENTER)

add_text(slide1, Inches(2), Inches(4.3), Inches(9.333), Inches(0.5),
         "FHIR R4 Compliant  ·  Multi-Agent NLP-to-SQL  ·  Human-in-the-Loop Governance",
         font_size=14, color=GRAY_400, alignment=PP_ALIGN.CENTER)

add_text(slide1, Inches(2), Inches(5.5), Inches(9.333), Inches(0.4),
         "React  ·  Django  ·  FastAPI  ·  LangGraph  ·  MCP  ·  PostgreSQL  ·  Redis  ·  Socket.IO",
         font_size=12, color=GRAY_500, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — System Overview
# ════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

add_text(slide2, Inches(0.5), Inches(0.3), Inches(5), Inches(0.5),
         "System Overview", font_size=28, color=WHITE, bold=True)
add_text(slide2, Inches(0.5), Inches(0.75), Inches(10), Inches(0.4),
         "6 microservices orchestrated via Docker Compose, unified by a single PostgreSQL database and Redis instance",
         font_size=12, color=GRAY_400)

# Service cards — row 1
services_row1 = [
    ("Frontend", "React + Vite + Nginx", "Port 3040", BLUE, BLUE_BG,
     ["Chat NLP Interface", "FHIR Data Browser (12 resources)", "Dashboard & Statistics", "HITL Approval Panel", "Agent Monitor & Observability"]),
    ("WebSocket Server", "Node.js + Socket.IO", "Port 3002", YELLOW, YELLOW_BG,
     ["Real-time bidirectional events", "submit_query / submit_decision", "approval_required / query_result", "Redis Pub/Sub bridge", "Session management"]),
    ("Agent Orchestrator", "FastAPI + LangGraph", "Port 8001", PINK, PINK_BG,
     ["9-node LangGraph workflow", "SQLAgent (GPT-4o-mini)", "ClassifierAgent (PHI + Risk)", "ExecutorAgent (asyncpg)", "HITLAgent (interrupt/resume)"]),
]

for i, (title, tech, port, color, bg, items) in enumerate(services_row1):
    left = Inches(0.5) + Inches(i * 4.2)
    card = add_shape(slide2, left, Inches(1.3), Inches(3.9), Inches(2.6), bg, color, Pt(2))
    add_text(slide2, left + Inches(0.25), Inches(1.4), Inches(3.4), Inches(0.3), title, 15, color, True)
    add_text(slide2, left + Inches(0.25), Inches(1.7), Inches(3.4), Inches(0.25), f"{tech}  ·  {port}", 10, GRAY_400)
    lines = [(f"  •  {item}", False, GRAY_300, 10) for item in items]
    add_rich_text(slide2, left + Inches(0.25), Inches(2.05), Inches(3.4), Inches(1.7), lines)

# Service cards — row 2
services_row2 = [
    ("MCP Server", "Node.js + Express + pg", "Port 3001", GREEN, GREEN_BG,
     ["GET /schema (fhir_* only)", "POST /validate-sql", "GET /tools & POST /tools/:name", "GET /resources (FHIR bundles)", "SQL injection detection"]),
    ("Django Backend", "Gunicorn + DRF + ORM", "Port 8000", TEAL, TEAL_BG,
     ["/api/fhir/* — CRUD + Search", "/api/audit/* — Audit trail", "/api/hitl/* — Approval tasks", "/api/observability/* — Metrics", "/admin/ — Django Admin"]),
    ("Data Layer", "PostgreSQL 16 + Redis 7", "Ports 5432 / 6379", INDIGO, INDIGO_BG,
     ["12 FHIR R4 tables", "audit_log + agent_interactions", "Redis DB 0-3 partitioning", "Schema, Cache, State, Pub/Sub", "LangGraph checkpoints"]),
]

for i, (title, tech, port, color, bg, items) in enumerate(services_row2):
    left = Inches(0.5) + Inches(i * 4.2)
    card = add_shape(slide2, left, Inches(4.15), Inches(3.9), Inches(2.6), bg, color, Pt(2))
    add_text(slide2, left + Inches(0.25), Inches(4.25), Inches(3.4), Inches(0.3), title, 15, color, True)
    add_text(slide2, left + Inches(0.25), Inches(4.55), Inches(3.4), Inches(0.25), f"{tech}  ·  {port}", 10, GRAY_400)
    lines = [(f"  •  {item}", False, GRAY_300, 10) for item in items]
    add_rich_text(slide2, left + Inches(0.25), Inches(4.9), Inches(3.4), Inches(1.7), lines)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Request Flow Architecture
# ════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, BG_DARK)

add_text(slide3, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
         "Request Flow Architecture", font_size=28, color=WHITE, bold=True)
add_text(slide3, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
         "End-to-end flow from natural language query to database result with Human-in-the-Loop governance",
         font_size=12, color=GRAY_400)

flow_steps = [
    ("1", "User Query", "Chat UI sends\nnatural language\nvia Socket.IO", BLUE, BLUE_BG),
    ("2", "WS Server", "Receives event,\nPOSTs to Agent\nOrchestrator", YELLOW, YELLOW_BG),
    ("3", "Fetch Schema", "GET /schema from\nMCP Server\n(fhir_* tables)", GREEN, GREEN_BG),
    ("4", "Generate SQL", "SQLAgent +\nGPT-4o-mini\nFHIR R4 prompt", PINK, PINK_BG),
    ("5", "Classify & Guard", "PHI detection,\nrisk scoring,\nguardrail checks", PURPLE, PURPLE_BG),
    ("6", "Route", "READ → auto-exec\nWRITE → HITL gate\nUNSAFE → reject", YELLOW, YELLOW_BG),
    ("7", "Execute SQL", "ExecutorAgent\nruns query via\nasyncpg", GREEN, GREEN_BG),
    ("8", "Present", "LLM summarizes\nraw results into\nhuman text", PINK, PINK_BG),
    ("9", "Audit Log", "Record to\naudit_log &\nagent_interactions", INDIGO, INDIGO_BG),
]

for i, (num, title, desc, color, bg) in enumerate(flow_steps):
    left = Inches(0.35) + Inches(i * 1.42)
    top = Inches(1.4)
    w = Inches(1.3)
    h = Inches(1.6)
    add_shape(slide3, left, top, w, h, bg, color, Pt(2))
    # Number circle
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.45), top + Inches(0.08), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide3, left + Inches(0.08), top + Inches(0.55), w - Inches(0.16), Inches(0.25), title, 11, color, True, PP_ALIGN.CENTER)
    add_text(slide3, left + Inches(0.08), top + Inches(0.8), w - Inches(0.16), Inches(0.7), desc, 9, GRAY_300, False, PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(flow_steps) - 1:
        arrow_left = left + w + Inches(0.01)
        arrow_top = top + h / 2
        arr = slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top - Inches(0.08), Inches(0.12), Inches(0.16))
        arr.fill.solid()
        arr.fill.fore_color.rgb = GRAY_500
        arr.line.fill.background()

# HITL approval box
add_shape(slide3, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.6), BG_CARD, YELLOW, Pt(2))
add_text(slide3, Inches(0.8), Inches(3.5), Inches(4), Inches(0.35), "Human-in-the-Loop (HITL) Approval Flow", 14, YELLOW, True)

hitl_items = [
    ("WRITE queries or PHI access detected", False, GRAY_300, 10),
    ("LangGraph interrupt() pauses workflow at hitl_gate node", False, GRAY_300, 10),
    ("WebSocket emits approval_required → Approval Panel appears in Chat UI", False, GRAY_300, 10),
    ("Reviewer sees: Generated SQL, Risk Score (%), Risk Assessment, Query Type", False, GRAY_300, 10),
    ("APPROVED → resume workflow → execute_sql → present_results → audit", False, GREEN, 10),
    ("REJECTED → handle_rejection → log_audit → notify user", False, RED, 10),
]
add_rich_text(slide3, Inches(0.8), Inches(3.85), Inches(5.5), Inches(1.1),
              [(f"  •  {t}", b, c, s) for t, b, c, s in hitl_items])

# Toxicity & Safety box
add_shape(slide3, Inches(6.8), Inches(3.85), Inches(5.8), Inches(1.0), BG_CARD_LIGHT, RED, Pt(1))
add_text(slide3, Inches(7.0), Inches(3.9), Inches(3), Inches(0.3), "Security & Safety Layers", 12, RED, True)
safety_items = [
    ("Toxicity filter — blocks hate speech, discrimination, malicious intent", False, GRAY_300, 9),
    ("PHI detector — identifies 40+ sensitive column patterns (HIPAA)", False, GRAY_300, 9),
    ("Blocked columns — SSN, Passport, Drivers License always denied", False, GRAY_300, 9),
    ("SQL injection detection — ;DROP, UNION SELECT, 'OR 1=1' patterns", False, GRAY_300, 9),
    ("SELECT * prohibition — enforces data minimization principle", False, GRAY_300, 9),
    ("Audit logging — every query, decision, and result recorded", False, GRAY_300, 9),
]
add_rich_text(slide3, Inches(7.0), Inches(4.2), Inches(5.5), Inches(0.65),
              [(f"  •  {t}", b, c, s) for t, b, c, s in safety_items])

# Data flow summary
add_shape(slide3, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), BG_CARD, BLUE, Pt(1))
add_text(slide3, Inches(0.8), Inches(5.4), Inches(4), Inches(0.3), "Two Data Access Paths", 14, BLUE, True)

add_rich_text(slide3, Inches(0.8), Inches(5.75), Inches(5.8), Inches(1.2), [
    ("Chat NLP Path (Agent-mediated):", True, BLUE_LIGHT, 11),
    ("  Frontend → Socket.IO → WS Server → Agent Orchestrator", False, GRAY_300, 10),
    ("  → MCP Server (schema) → SQLAgent (generate) → Classifier (guard)", False, GRAY_300, 10),
    ("  → Executor (asyncpg → PostgreSQL) → LLM Summary → Frontend", False, GRAY_300, 10),
])

add_rich_text(slide3, Inches(6.8), Inches(5.75), Inches(5.8), Inches(1.2), [
    ("FHIR Browser Path (Direct REST):", True, GREEN, 11),
    ("  Frontend → HTTP GET/POST → Django Backend (DRF)", False, GRAY_300, 10),
    ("  → Django ORM → PostgreSQL (fhir_* tables)", False, GRAY_300, 10),
    ("  → JSON response → React table/detail views", False, GRAY_300, 10),
])


# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — LangGraph Multi-Agent Workflow
# ════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, BG_DARK)

add_text(slide4, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
         "LangGraph Multi-Agent Workflow", font_size=28, color=WHITE, bold=True)
add_text(slide4, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
         "Stateful graph with conditional routing, interrupt/resume for HITL, and comprehensive audit logging",
         font_size=12, color=GRAY_400)

# Workflow nodes
nodes = [
    ("START", "", BG_CARD_LIGHT, WHITE, 0.5, 1.5, 1.0, 0.5),
    ("1. check_toxicity", "Filter harmful queries\nHate speech, discrimination,\nmalicious intent detection", PINK_BG, PINK, 1.8, 1.4, 2.4, 0.7),
    ("2. fetch_schema", "GET /schema from MCP\nReturns fhir_* tables\nCached in Redis DB 1", GREEN_BG, GREEN, 1.8, 2.3, 2.4, 0.7),
    ("3. generate_sql", "SQLAgent + GPT-4o-mini\nFHIR R4 prompt (12 tables)\nAGE() function guidance", PINK_BG, PINK, 1.8, 3.2, 2.4, 0.7),
    ("4. classify_query", "ClassifierAgent\nPHI column detection\nRisk score (0.0 – 1.0)", PURPLE_BG, PURPLE, 1.8, 4.1, 2.4, 0.7),
    ("5. check_guardrails", "SQL injection patterns\nSELECT * prohibition\nBlocked columns check", RED_BG, RED, 1.8, 5.0, 2.4, 0.7),
]

for label, desc, bg, color, left, top, w, h in nodes:
    add_shape(slide4, Inches(left), Inches(top), Inches(w), Inches(h), bg, color, Pt(2))
    add_text(slide4, Inches(left + 0.15), Inches(top + 0.05), Inches(w - 0.3), Inches(0.25), label, 11, color, True)
    if desc:
        add_text(slide4, Inches(left + 0.15), Inches(top + 0.28), Inches(w - 0.3), Inches(h - 0.3), desc, 8, GRAY_300)

# Routing diamonds
routes = [
    ("READ\n(Auto)", GREEN, GREEN_BG, 4.8, 5.1),
    ("WRITE\n(HITL)", YELLOW, YELLOW_BG, 6.3, 5.1),
    ("UNSAFE\n(Block)", RED, RED_BG, 7.8, 5.1),
]
for label, color, bg, left, top in routes:
    diamond = slide4.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(left), Inches(top), Inches(1.2), Inches(0.8))
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = bg
    diamond.line.color.rgb = color
    diamond.line.width = Pt(2)
    tf = diamond.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Right column nodes
right_nodes = [
    ("6. hitl_gate", "LangGraph interrupt()\nAwaits human decision\nApprove / Reject + Notes", YELLOW_BG, YELLOW, 5.0, 3.2, 2.5, 0.7),
    ("7. execute_sql", "ExecutorAgent\nasyncpg direct to PostgreSQL\nParameterized queries", GREEN_BG, GREEN, 5.0, 1.4, 2.5, 0.7),
    ("8. present_results", "LLM summarization\nRaw SQL → human text\nMarkdown formatting", PINK_BG, PINK, 8.0, 1.4, 2.5, 0.7),
    ("9. log_audit", "audit_log table\nagent_interactions table\nFull query lifecycle", INDIGO_BG, INDIGO, 8.0, 2.3, 2.5, 0.7),
    ("handle_rejection", "Block message\nViolation details\nLog and notify user", RED_BG, RED, 8.0, 3.2, 2.5, 0.7),
]
for label, desc, bg, color, left, top, w, h in right_nodes:
    add_shape(slide4, Inches(left), Inches(top), Inches(w), Inches(h), bg, color, Pt(2))
    add_text(slide4, Inches(left + 0.15), Inches(top + 0.05), Inches(w - 0.3), Inches(0.25), label, 11, color, True)
    add_text(slide4, Inches(left + 0.15), Inches(top + 0.28), Inches(w - 0.3), Inches(h - 0.3), desc, 8, GRAY_300)

# END node
end = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(2.5), Inches(0.8), Inches(0.5))
end.fill.solid()
end.fill.fore_color.rgb = WHITE
end.line.fill.background()
tf = end.text_frame
p = tf.paragraphs[0]
p.text = "END"
p.font.size = Pt(12)
p.font.color.rgb = BG_DARK
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Agent info cards at bottom
add_shape(slide4, Inches(5.0), Inches(4.2), Inches(7.5), Inches(2.5), BG_CARD, GRAY_500, Pt(1))
add_text(slide4, Inches(5.2), Inches(4.3), Inches(4), Inches(0.3), "Agent Details", 14, WHITE, True)

agents_info = [
    ("SQLAgent", "GPT-4o-mini LLM, FHIR R4 prompt with all 12 table schemas,\nrelationships, AGE() guidance, SELECT * prohibition", PINK),
    ("ClassifierAgent", "Rule-based PHI detection (40+ patterns), risk scoring,\ntoxicity analysis, READ/WRITE/UNSAFE classification", PURPLE),
    ("ExecutorAgent", "asyncpg connection pool, parameterized query execution,\nPHI masking on results, timeout protection", GREEN),
    ("HITLAgent", "LangGraph interrupt/resume, Redis pub/sub notifications,\napproval task management, reviewer audit trail", YELLOW),
]

for i, (name, desc, color) in enumerate(agents_info):
    y = Inches(4.65) + Inches(i * 0.48)
    add_shape(slide4, Inches(5.2), y, Inches(1.6), Inches(0.35), BG_CARD_LIGHT, color, Pt(1))
    add_text(slide4, Inches(5.3), y + Inches(0.03), Inches(1.4), Inches(0.3), name, 10, color, True, PP_ALIGN.CENTER)
    add_text(slide4, Inches(6.9), y + Inches(0.02), Inches(5.4), Inches(0.35), desc, 9, GRAY_300)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — FHIR R4 Database Schema
# ════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, BG_DARK)

add_text(slide5, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
         "FHIR R4 Database Schema", font_size=28, color=WHITE, bold=True)
add_text(slide5, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
         "12 FHIR R4 compliant resources in PostgreSQL with proper foreign key relationships",
         font_size=12, color=GRAY_400)

# Resource cards
resources = [
    ("fhir_patients", "Patient", "name_given, name_family, gender,\nbirth_date, address_*, telecom_*", BLUE, BLUE_BG),
    ("fhir_practitioners", "Practitioner", "name_given, name_family, NPI,\nqualification, specialty", TEAL, TEAL_BG),
    ("fhir_organizations", "Organization", "name, type, address,\ntelecom_phone, telecom_email", GREEN, GREEN_BG),
    ("fhir_encounters", "Encounter", "status, class, type, period_start,\nsubject_id → patients", YELLOW, YELLOW_BG),
    ("fhir_conditions", "Condition", "code (ICD-10), clinical_status,\nonset, subject_id → patients", PINK, PINK_BG),
    ("fhir_observations", "Observation", "code (LOINC), value, unit,\ncategory, interpretation", PURPLE, PURPLE_BG),
    ("fhir_medication_requests", "MedicationRequest", "medication (RxNorm), status,\ndosage, subject_id → patients", RED, RED_BG),
    ("fhir_allergy_intolerances", "AllergyIntolerance", "code (SNOMED), criticality,\ncategory, patient_id → patients", YELLOW, YELLOW_BG),
    ("fhir_procedures", "Procedure", "code (SNOMED/CPT), status,\nperformed_date, body_site", GREEN, GREEN_BG),
    ("fhir_immunizations", "Immunization", "vaccine (CVX), lot_number,\noccurrence, patient_id → patients", TEAL, TEAL_BG),
    ("fhir_diagnostic_reports", "DiagnosticReport", "code (LOINC), conclusion,\neffective_date, status", BLUE, BLUE_BG),
    ("fhir_care_plans", "CarePlan", "title, status, intent,\nperiod, description", INDIGO, INDIGO_BG),
]

for i, (table, resource, desc, color, bg) in enumerate(resources):
    col = i % 4
    row = i // 4
    left = Inches(0.5) + Inches(col * 3.15)
    top = Inches(1.3) + Inches(row * 1.55)
    w = Inches(2.95)
    h = Inches(1.35)
    add_shape(slide5, left, top, w, h, bg, color, Pt(2))
    add_text(slide5, left + Inches(0.15), top + Inches(0.08), w - Inches(0.3), Inches(0.22), resource, 12, color, True)
    add_text(slide5, left + Inches(0.15), top + Inches(0.32), w - Inches(0.3), Inches(0.2), table, 8, GRAY_400)
    add_text(slide5, left + Inches(0.15), top + Inches(0.55), w - Inches(0.3), Inches(0.75), desc, 9, GRAY_300)

# Relationships box
add_shape(slide5, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.2), BG_CARD, INDIGO, Pt(1))
add_text(slide5, Inches(0.8), Inches(6.0), Inches(4), Inches(0.3), "Key Relationships & Coding Systems", 13, INDIGO, True)

add_rich_text(slide5, Inches(0.8), Inches(6.3), Inches(5.5), Inches(0.8), [
    ("FK: subject_id → fhir_patients (Encounter, Condition, Observation,", False, GRAY_300, 9),
    ("    MedicationRequest, Procedure, DiagnosticReport, CarePlan)", False, GRAY_300, 9),
    ("FK: patient_id → fhir_patients (AllergyIntolerance, Immunization)", False, YELLOW, 9),
    ("FK: encounter_id → fhir_encounters (all clinical resources)", False, GRAY_300, 9),
])

add_rich_text(slide5, Inches(6.8), Inches(6.3), Inches(5.5), Inches(0.8), [
    ("ICD-10-CM — Diagnoses (fhir_conditions.code_code)", False, GRAY_300, 9),
    ("SNOMED CT — Clinical findings & procedures", False, GRAY_300, 9),
    ("LOINC — Laboratory & vital signs observations", False, GRAY_300, 9),
    ("RxNorm — Medications  ·  CVX — Vaccines  ·  CPT — Procedures", False, GRAY_300, 9),
])


# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — MCP Server & Tools
# ════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, BG_DARK)

add_text(slide6, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
         "MCP Server & Tool Architecture", font_size=28, color=WHITE, bold=True)
add_text(slide6, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
         "Model Context Protocol server provides schema, validation, and tool execution for the agent orchestrator",
         font_size=12, color=GRAY_400)

# MCP Endpoints
endpoints = [
    ("GET /schema", "Returns database schema for fhir_* tables only.\nFiltered via: WHERE table_name LIKE 'fhir_%'\nIncludes column names, types, nullability.", GREEN),
    ("POST /validate-sql", "Validates SQL before execution.\nBlocks: DROP, TRUNCATE, ALTER, GRANT, REVOKE\nDetects injection: ;DROP, UNION SELECT, 'OR 1=1'\nWarns on missing LIMIT clause.", RED),
    ("GET /tools", "Lists all available MCP tools:\nget_fhir_patient, search_fhir_patients,\nsearch_fhir_conditions, search_fhir_observations,\nsearch_fhir_medications, get_fhir_statistics,\nlist_tables, get_table_schema", BLUE),
    ("POST /tools/:name", "Executes a specific tool with arguments.\nPasses pg.Pool connection to tool handler.\nReturns structured JSON results.\nUsed for targeted data access.", PURPLE),
    ("GET /resources", "Lists FHIR resource bundles:\nPatient Bundle, Practitioner Bundle,\nOrganization Bundle, Encounter Bundle,\nFHIR relationship mappings.", YELLOW),
    ("GET /health", "Service health check.\nVerifies PostgreSQL connection.\nVerifies Redis connection.\nReturns service status.", GRAY_400),
]

for i, (ep, desc, color) in enumerate(endpoints):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + Inches(col * 4.2)
    top = Inches(1.3) + Inches(row * 2.5)
    w = Inches(3.9)
    h = Inches(2.2)
    bg = BG_CARD
    add_shape(slide6, left, top, w, h, bg, color, Pt(2))
    add_text(slide6, left + Inches(0.2), top + Inches(0.12), w - Inches(0.4), Inches(0.3), ep, 14, color, True)
    add_text(slide6, left + Inches(0.2), top + Inches(0.45), w - Inches(0.4), Inches(1.6), desc, 10, GRAY_300)

# Connection info
add_shape(slide6, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.9), BG_CARD, GREEN, Pt(1))
add_text(slide6, Inches(0.8), Inches(6.35), Inches(3), Inches(0.3), "Connection Architecture", 13, GREEN, True)
add_rich_text(slide6, Inches(0.8), Inches(6.65), Inches(11), Inches(0.5), [
    ("Agent Orchestrator → HTTP → MCP Server (Port 3001) → pg.Pool → PostgreSQL (Port 5432) → fhir_* tables", False, GRAY_300, 10),
    ("MCP Server → Redis (DB 1) for schema caching  ·  Schema refreshed on each agent request  ·  SQL validation is stateless", False, GRAY_400, 9),
])


# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Security & Governance
# ════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, BG_DARK)

add_text(slide7, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
         "Security, PHI & Governance", font_size=28, color=WHITE, bold=True)
add_text(slide7, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
         "Multi-layered HIPAA-aligned protection with automated and human governance controls",
         font_size=12, color=GRAY_400)

sec_cards = [
    ("PHI Detection Engine", RED, RED_BG, [
        "40+ sensitive column patterns detected",
        "Blocked columns: SSN, Passport, Drivers License",
        "HITL required: Name, Birth Date, Phone, Email, MRN",
        "Supports both Synthea & FHIR column naming",
        "Masking levels: FULL, PARTIAL, REDACT, NONE",
    ]),
    ("Toxicity Filter", PINK, PINK_BG, [
        "Hate speech & discrimination detection",
        "Data exfiltration attempt blocking",
        "Malicious intent recognition",
        "Privacy violation detection",
        "Bulk sensitive data request blocking",
    ]),
    ("SQL Safety Guards", YELLOW, YELLOW_BG, [
        "SELECT * always blocked (data minimization)",
        "DROP/TRUNCATE/ALTER/GRANT blocked",
        "SQL injection pattern detection",
        "UPDATE/DELETE without WHERE blocked",
        "LIMIT enforcement on all queries",
    ]),
    ("Human-in-the-Loop", BLUE, BLUE_BG, [
        "All WRITE operations require approval",
        "PHI column access requires approval",
        "Risk score visualization (0-100%)",
        "Reviewer notes and audit trail",
        "LangGraph interrupt/resume mechanism",
    ]),
    ("Audit Trail", INDIGO, INDIGO_BG, [
        "Every query logged: NL → SQL → result",
        "Classification & risk scores recorded",
        "Reviewer decisions with timestamps",
        "Agent-to-agent interaction logs",
        "Full session reconstruction capability",
    ]),
    ("Data Classification", GREEN, GREEN_BG, [
        "READ — auto-executed, low risk",
        "WRITE — requires HITL approval",
        "UNSAFE — immediately blocked",
        "Risk score: 0.0 (safe) to 1.0 (blocked)",
        "Dynamic scoring based on PHI + complexity",
    ]),
]

for i, (title, color, bg, items) in enumerate(sec_cards):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + Inches(col * 4.2)
    top = Inches(1.3) + Inches(row * 2.7)
    w = Inches(3.9)
    h = Inches(2.45)
    add_shape(slide7, left, top, w, h, bg, color, Pt(2))
    add_text(slide7, left + Inches(0.2), top + Inches(0.12), w - Inches(0.4), Inches(0.3), title, 14, color, True)
    lines = [(f"  •  {item}", False, GRAY_300, 10) for item in items]
    add_rich_text(slide7, left + Inches(0.2), top + Inches(0.45), w - Inches(0.4), Inches(1.8), lines)


# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Technology Stack
# ════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, BG_DARK)

add_text(slide8, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
         "Technology Stack", font_size=28, color=WHITE, bold=True)

stack_sections = [
    ("Frontend", BLUE, BLUE_BG, [
        ("React 18", "Component-based UI framework"),
        ("TypeScript", "Type-safe development"),
        ("Vite", "Fast build tool & HMR"),
        ("Tailwind CSS", "Utility-first styling"),
        ("Zustand", "Lightweight state management"),
        ("React Router v6", "Client-side routing"),
        ("Socket.IO Client", "WebSocket communication"),
        ("Axios", "HTTP REST API client"),
        ("Lucide React", "Icon library"),
    ]),
    ("Backend", GREEN, GREEN_BG, [
        ("Django 4.2", "Web framework + ORM"),
        ("Django REST Framework", "REST API toolkit"),
        ("Gunicorn", "WSGI HTTP server"),
        ("django-filters", "QuerySet filtering"),
        ("django-cors-headers", "CORS middleware"),
        ("django-redis", "Cache backend"),
        ("dj-database-url", "Database URL parsing"),
        ("psycopg2", "PostgreSQL adapter"),
        ("", ""),
    ]),
    ("AI / Agents", PINK, PINK_BG, [
        ("LangGraph", "Stateful agent workflow graph"),
        ("LangChain", "LLM orchestration framework"),
        ("OpenAI GPT-4o-mini", "SQL generation LLM"),
        ("FastAPI", "Async API framework"),
        ("Uvicorn", "ASGI server"),
        ("asyncpg", "Async PostgreSQL driver"),
        ("LangSmith / Langfuse", "Observability & tracing"),
        ("httpx", "Async HTTP client"),
        ("", ""),
    ]),
    ("Infrastructure", INDIGO, INDIGO_BG, [
        ("Docker Compose", "Multi-service orchestration"),
        ("PostgreSQL 16", "Primary database (Alpine)"),
        ("Redis 7", "Cache + Pub/Sub (Alpine)"),
        ("Nginx", "Static file server / reverse proxy"),
        ("Node.js 20", "MCP + WS server runtime"),
        ("Express.js", "MCP HTTP framework"),
        ("Socket.IO", "Real-time WebSocket server"),
        ("node-postgres (pg)", "PostgreSQL connection pool"),
        ("", ""),
    ]),
]

for i, (title, color, bg, items) in enumerate(stack_sections):
    left = Inches(0.5) + Inches(i * 3.15)
    w = Inches(2.95)
    add_shape(slide8, left, Inches(1.0), w, Inches(6.0), bg, color, Pt(2))
    add_text(slide8, left + Inches(0.15), Inches(1.1), w - Inches(0.3), Inches(0.3), title, 15, color, True, PP_ALIGN.CENTER)

    for j, (tech, desc) in enumerate(items):
        if not tech:
            continue
        y = Inches(1.5) + Inches(j * 0.55)
        add_shape(slide8, left + Inches(0.1), y, w - Inches(0.2), Inches(0.45), BG_CARD_LIGHT)
        add_text(slide8, left + Inches(0.2), y + Inches(0.02), w - Inches(0.4), Inches(0.2), tech, 10, color, True)
        add_text(slide8, left + Inches(0.2), y + Inches(0.22), w - Inches(0.4), Inches(0.2), desc, 8, GRAY_400)


# ════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════
output_path = "/home/user/Health_Assistant/docs/Healthcare_Platform_Architecture.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
